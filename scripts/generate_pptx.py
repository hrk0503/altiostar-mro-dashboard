"""Generates a PowerPoint presentation (docs/presentation.pptx) summarizing the MRO training pipeline project.
"""
from __future__ import annotations

import sys
import subprocess
from pathlib import Path

# Programmatically check and install python-pptx if missing
try:
    import pptx
except ImportError:
    print("python-pptx is not installed. Installing it via pip...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_presentation() -> None:
    prs = Presentation()
    
    # Define color scheme
    color_bg = RGBColor(11, 15, 25)        # Dark Blue
    color_card = RGBColor(17, 24, 39)       # Card BG
    color_text = RGBColor(243, 244, 246)    # Off-white
    color_sub = RGBColor(156, 163, 175)     # Gray text
    color_primary = RGBColor(56, 189, 248)  # Light Blue Accent
    color_accent = RGBColor(129, 140, 248)  # Indigo
    
    # Slide Dimensions (16:9 widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Background helper
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_bg
        
    # Helper to add standard header to slide
    def add_slide_header(slide, title_text, tag_text=""):
        # Header text box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color_text
        
        if tag_text:
            p_tag = tf.add_paragraph()
            p_tag.text = f"// {tag_text}"
            p_tag.font.name = "Arial"
            p_tag.font.size = Pt(13)
            p_tag.font.bold = True
            p_tag.font.color.rgb = color_primary
            
    # Helper to add bullet point list
    def add_bullet_list(slide, items, left, top, width, height, font_size=16):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        for idx, item in enumerate(items):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = "• " + item[0]
            p.font.name = "Arial"
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = color_text
            p.space_after = Pt(2)
            
            p_desc = tf.add_paragraph()
            p_desc.text = "  " + item[1]
            p_desc.font.name = "Arial"
            p_desc.font.size = Pt(font_size - 3)
            p_desc.font.color.rgb = color_sub
            p_desc.space_after = Pt(12)
            
        return box

    # Slide 1: Title Slide (Blank Layout)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "5G MRO RL Pipeline Deep-Dive"
    p.font.name = "Arial"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = color_text
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "Under-the-Hood Data Pipeline, Gymnasium Simulation, and Multi-Agent Architecture"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.color.rgb = color_primary
    p2.space_after = Pt(48)
    
    p3 = tf.add_paragraph()
    p3.text = "AltioStar Tokyo MRO Stream | Devika · Harshit · Ananyaa · Shourya"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = color_accent
    p3.space_after = Pt(4)
    
    p4 = tf.add_paragraph()
    p4.text = "WINNIIO × Amity 2026 | June 2026"
    p4.font.name = "Arial"
    p4.font.size = Pt(13)
    p4.font.color.rgb = color_sub

    # Slide 2: Ingestion & Schema Mapper
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "The Ingestion Data Pipeline", "Data Engineering")
    
    bullets_p1 = [
        ("Raw Operator CSV Files", "Ingests 4 base files: site database (75 cells), neighbor relations (763 pairs), PM data (216K rows), and monthly KPI summary (75 rows)."),
        ("Relation-Level Compilation", "Merges and compiles cell-level PM and relation parameters into a high-granularity 2.2M-row relation-level PM dataset (pm_data_relation_level.csv)."),
        ("Set-Based Schema Mapper", "Auto-detects CSV types and maps columns. Uses set signature matching (cols.issubset) in schema_mapper.py to remain 100% order-independent (shuffled column proof)."),
        ("Double Caching Performance", "Caches dataframes and pre-grouped relation slices in memory to reduce env boot time from minutes to under 50 milliseconds.")
    ]
    add_bullet_list(slide, bullets_p1, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
    ctf = card_box.text_frame
    ctf.word_wrap = True
    
    cp = ctf.paragraphs[0]
    cp.text = "📊 CSV to State Lifecycle"
    cp.font.name = "Arial"
    cp.font.size = Pt(22)
    cp.font.bold = True
    cp.font.color.rgb = color_primary
    cp.space_after = Pt(16)
    
    cp2 = ctf.add_paragraph()
    cp2.text = "1. Raw CSV Ingestion\n   ↓\n2. Set-based Signature Mapping (schema_mapper.py)\n   ↓\n3. Pre-Grouped Relation Indices Caching\n   ↓\n4. Fast O(1) Gymnasium Env Access"
    cp2.font.name = "Arial"
    cp2.font.size = Pt(16)
    cp2.font.color.rgb = color_text

    # Slide 3: Gymnasium MRO Env
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Gymnasium Environment Simulation", "Under the Hood")
    
    bullets_p2 = [
        ("Observation Space", "9-dimensional vector per relation (handover attempts, successes, failures: too early, too late, wrong cell, correct cell, ping-pongs, average RSRP/SINR)."),
        ("Action Space", "Continuous Cell Individual Offset (CIO) delta adjustments per neighbor relation, bounded strictly to [-2.0, 2.0] dB."),
        ("Data-Driven State Transitions", "No slow physical propagation math. Environment performs binary search to find the nearest historical CIO database configuration and samples outcomes from the real distribution."),
        ("Fast Step Simulation", "Optimized grouping and caches yield O(1) transitions, allowing 100k rollout steps to run in seconds on CPU.")
    ]
    add_bullet_list(slide, bullets_p2, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
    ctf = card_box.text_frame
    ctf.word_wrap = True
    
    cp = ctf.paragraphs[0]
    cp.text = "🧠 Simulation Model"
    cp.font.name = "Arial"
    cp.font.size = Pt(22)
    cp.font.bold = True
    cp.font.color.rgb = color_primary
    cp.space_after = Pt(16)
    
    cp2 = ctf.add_paragraph()
    cp2.text = "Physics modeling is slow and error-prone. By slicing the 2.2M-row database and caching outcomes per relation, the env acts as a high-fidelity lookup simulator that behaves exactly like the real radio network."
    cp2.font.name = "Arial"
    cp2.font.size = Pt(16)
    cp2.font.color.rgb = color_text

    # Slide 4: Reward Tuning & Ship Gate
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Reward Engineering & Ship Gate", "Model Safety")
    
    bullets_p3 = [
        ("v2 Rate-Based Reward", "Normalizes by percentage rates directly (bounded and consistent): Reward = HOSR * 1.0 - failure_rate * 5.0 - PP_rate * 2.0."),
        ("Asymmetric Boundary Penalties", "Introduces asymmetric barrier functions to env (punishes HOSR < 95% and PP > 5% with large negative penalties) to guide PPO policy convergence."),
        ("Automated Ship Gate Validator", "Built ship_gate.py checking results against strictly HOSR > 99% and PP < 5%. Outputs exit code 0/1 for CI/CD gates."),
        ("Statistical Significance Check", "Evaluates sweep results across 5 independent seeds to ensure repeatability and prevent random seed selection bias.")
    ]
    add_bullet_list(slide, bullets_p3, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
    ctf = card_box.text_frame
    ctf.word_wrap = True
    
    cp = ctf.paragraphs[0]
    cp.text = "🛡️ Ship Gate Criteria"
    cp.font.name = "Arial"
    cp.font.size = Pt(22)
    cp.font.bold = True
    cp.font.color.rgb = color_primary
    cp.space_after = Pt(16)
    
    cp2 = ctf.add_paragraph()
    cp2.text = "• HOSR: Strictly > 99.0% (99.0% = FAIL)\n• Ping-Pong: Strictly < 5.0% (5.0% = FAIL)\n\nChecked automatically on every sweep results JSON file before code merges to staging."
    cp2.font.name = "Arial"
    cp2.font.size = Pt(16)
    cp2.font.color.rgb = color_text

    # Slide 5: Agent Roles
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Multi-Agent System Architecture", "Agent Execution Roles")
    
    roles_left = [
        ("Atlas (Product Owner)", "Orchestrates sprint goals, phase gates, daily reports, and blocker escalations."),
        ("Pipeline (Data Engineer)", "Responsible for CSV ingestion, schema mapper auto-detection, data validation, and Parquet exports."),
        ("Spectrum (RF Domain Lead)", "RF domain expert. Enforces 3GPP constraints, configures valid CIO ranges, and guides reward calibration."),
        ("Forge (ML Engineer)", "Builds the Gymnasium environment, handles PPO convergence loops, and runs Optuna parameter sweeps.")
    ]
    
    roles_right = [
        ("Deploy (DevOps)", "Docker containment, environment locks, and K8s manifests targeting CU-CP K8s pods."),
        ("Lens (Visualization)", "Generates interactive Streamlit dashboards, HTML presentations, and PowerPoint decks."),
        ("Sentinel (QA & Validation)", "Property-based tests (Hypothesis), adversarial input audits, and automated ship gate verification.")
    ]
    
    add_bullet_list(slide, roles_left, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), font_size=15)
    add_bullet_list(slide, roles_right, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0), font_size=15)

    # Slide 6: O-RAN MRO Technical Architecture (Screenshot)
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    img_arch_path = Path("docs/screenshots/architecture.png")
    if img_arch_path.exists():
        slide.shapes.add_picture(str(img_arch_path), Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    else:
        add_slide_header(slide, "O-RAN MRO Technical Architecture", "System Design Overview")
        card_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
        ctf = card_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = "🏗️ End-to-End System Design"
        cp.font.name = "Arial"
        cp.font.size = Pt(22)
        cp.font.bold = True
        cp.font.color.rgb = color_primary
        cp.space_after = Pt(16)
        cp2 = ctf.add_paragraph()
        cp2.text = "Visualizes the end-to-end O-RAN MRO pipeline, from raw CSV ingestion to ONNX model export for K8s deployment."
        cp2.font.name = "Arial"
        cp2.font.size = Pt(16)
        cp2.font.color.rgb = color_text

    # Slide 7: Technical Architecture Breakdown (Explanation)
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Technical Pipeline Architecture", "System Design Breakdown")
    
    bullets_arch = [
        ("Raw Ingestion & Schema Mapper", "Reads PM Counters, Site Database, and Neighbor Relations CSVs. Maps columns dynamically using order-independent set signature matching (cols.issubset)."),
        ("In-Memory Double Cache Slices", "Groups compiled 2.2M-row relation-level PM records and caches indices in memory, reducing env boot and step lookup times to under 50 microseconds."),
        ("Gymnasium RL Simulator", "Models network transitions using nearest-CIO lookup. Exposes a 9D Observation Space and a continuous Action Space [-2.0, 2.0] dB delta for CIO offsets."),
        ("PPO Sweep & Ship Gate Validation", "Trains Stable-Baselines3 PPO variants (v0-v3) with Optuna sweeps. Enforces strict ship gate criteria (HOSR > 99.0%, Ping-Pong < 5.0%) before exporting ONNX with output clamping.")
    ]
    add_bullet_list(slide, bullets_arch, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
    ctf = card_box.text_frame
    ctf.word_wrap = True
    
    cp = ctf.paragraphs[0]
    cp.text = "🔌 O-RAN Deployment Flow"
    cp.font.name = "Arial"
    cp.font.size = Pt(22)
    cp.font.bold = True
    cp.font.color.rgb = color_primary
    cp.space_after = Pt(16)
    
    cp2 = ctf.add_paragraph()
    cp2.text = "The pipeline is built as a closed loop. Historical metrics feed the Simulator, PPO policy optimizes CIO offsets, validation checks quality gates, and ONNX models compile directly into K8s pods at baseband."
    cp2.font.name = "Arial"
    cp2.font.size = Pt(16)
    cp2.font.color.rgb = color_text

    # Slide 8: Streamlit Management Dashboard - Overview & Health
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Streamlit Management Dashboard", "Interactive Overview & Health")
    
    bullets_dash = [
        ("Executive Performance KPI Cards", "Displays live metrics including Handover Success Rate (HOSR), Ping-Pong rate, PRB usage, and active UE count."),
        ("Cell Health Distribution (Donut)", "Classifies the 75-cell cluster in Shibuya: 23 Healthy (green), 46 Warning (orange), and 6 Critical (red)."),
        ("Operational Health Breakdown", "Provides direct visual summaries of degraded cells to focus engineering efforts where handovers fail most."),
        ("SLA Performance Audits", "Tracks real-time KPIs and flags sites that fail to maintain the strict G4 target thresholds.")
    ]
    add_bullet_list(slide, bullets_dash, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    # Embed the cell health donut screenshot
    img_donut_path = Path("docs/screenshots/cell_health_donut.png")
    if img_donut_path.exists():
        slide.shapes.add_picture(str(img_donut_path), Inches(7.8), Inches(1.6), Inches(4.5), Inches(4.8))
    else:
        card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
        ctf = card_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = "📊 Donut Chart"
        cp.font.name = "Arial"
        cp.font.size = Pt(22)
        cp.font.bold = True
        cp.font.color.rgb = color_primary
        cp.space_after = Pt(16)
        cp2 = ctf.add_paragraph()
        cp2.text = "Visualizes cell health distribution (Healthy/Warning/Critical) to identify problematic sites at a glance."
        cp2.font.name = "Arial"
        cp2.font.size = Pt(16)
        cp2.font.color.rgb = color_text

    # Slide 7: Interactive Cell Map & Sector Details
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Geographic Cell Map & Site Inspector", "Interactive Topology")
    
    bullets_map = [
        ("Geographic Cluster Visualizer", "Renders the 75-cell Shibuya cluster, showing site sectors, frequency bands, and geographical coordinates."),
        ("Physical Attribute Mapping", "Displays sector frequency band, electrical/mechanical tilt, and elevation height for individual sites."),
        ("Handover Relation Overlay", "Draws lines connecting neighbor cells to expose spatial spacing, coverage boundaries, and overlapping zones."),
        ("Cell Detail Panel", "Provides granular reports for selected cells (e.g. RKSB-001-1), including active UEs, average SINR/RSRP, and neighbor table offsets.")
    ]
    add_bullet_list(slide, bullets_map, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    # Embed the cell map / details screenshot
    img_map_path = Path("docs/screenshots/cell_map.png")
    if img_map_path.exists():
        slide.shapes.add_picture(str(img_map_path), Inches(7.5), Inches(2.2), Inches(5.2), Inches(3.14))
    else:
        card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
        ctf = card_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = "🗺️ Topology Inspector"
        cp.font.name = "Arial"
        cp.font.size = Pt(22)
        cp.font.bold = True
        cp.font.color.rgb = color_primary
        cp.space_after = Pt(16)
        cp2 = ctf.add_paragraph()
        cp2.text = "Exposes cell sector geometry, coverage overlaps, and neighbor relationship offsets on a geographical layout."
        cp2.font.name = "Arial"
        cp2.font.size = Pt(16)
        cp2.font.color.rgb = color_text

    # Slide 8: RL Variant Comparison & Performance Heatmap
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Variant Comparison & Performance Heatmap", "Model Evaluation")
    
    bullets_variant = [
        ("HO Success Rate Comparison", "Compares all 4 variants (v0-v3) under the 4 scenarios (baseline, rain fade, rush hour, tower failure) as a bar chart."),
        ("SLA Benchmark Lines", "Shows a 99% Target threshold (red line) and the Random Baseline performance of 79.2% (gray line)."),
        ("Multi-Metric Radar Fingerprint", "Tracks balanced criteria (low failure, low ping-pong, low too-early, low wrong-cell) on a radar plot."),
        ("Scenario Heatmap Grid", "Visualizes variant success rate rates directly, validating 99.99% convergence under baseline and rain fade scenarios.")
    ]
    add_bullet_list(slide, bullets_variant, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    # Embed stacked variant bar chart and heatmap/radar screenshot
    img_bar_path = Path("docs/screenshots/experiment_results_bar.png")
    img_heat_path = Path("docs/screenshots/heatmap_radar.png")
    if img_bar_path.exists() and img_heat_path.exists():
        slide.shapes.add_picture(str(img_bar_path), Inches(7.5), Inches(1.6), Inches(5.0), Inches(2.0))
        slide.shapes.add_picture(str(img_heat_path), Inches(7.5), Inches(3.8), Inches(5.0), Inches(2.0))
    else:
        card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
        ctf = card_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = "📈 Variant Analytics"
        cp.font.name = "Arial"
        cp.font.size = Pt(22)
        cp.font.bold = True
        cp.font.color.rgb = color_primary
        cp.space_after = Pt(16)
        cp2 = ctf.add_paragraph()
        cp2.text = "Exposes variant success rates across baseline and stress scenarios to ensure safety and stability."
        cp2.font.name = "Arial"
        cp2.font.size = Pt(16)
        cp2.font.color.rgb = color_text

    # Slide 9: Real-Time RL Simulation & Parameters
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Real-Time RL Simulation", "Simulation Sandbox")
    
    bullets_sim = [
        ("Gymnasium Step Telemetry", "Provides real-time state visualization including Cluster HOSR (98.85% at step 15), cumulative reward (4,393), and degraded cell counts (5)."),
        ("Live Simulation Parameters", "Exposes active environment factors including UE Load (1.0x), RSRP Offset (0 dB), and Failure Multiplier (1.0x)."),
        ("Continuous CIO Tuning Monitor", "Displays real-time action adjustments (CIO +0 dB) computed by the RL agent for active neighbor relations."),
        ("Variant Leaderboard Track", "Compares active success rate predictions across variants (v0: 99.99%) to highlight the optimal operational policy.")
    ]
    add_bullet_list(slide, bullets_sim, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    # Embed the realtime simulation screenshot
    img_sim_path = Path("docs/screenshots/realtime_simulation.png")
    if img_sim_path.exists():
        slide.shapes.add_picture(str(img_sim_path), Inches(7.5), Inches(2.3), Inches(5.2), Inches(2.87))
    else:
        card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
        ctf = card_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = "🎮 Simulation Monitor"
        cp.font.name = "Arial"
        cp.font.size = Pt(22)
        cp.font.bold = True
        cp.font.color.rgb = color_primary
        cp.space_after = Pt(16)
        cp2 = ctf.add_paragraph()
        cp2.text = "Allows operators to trigger scenarios and visually verify agent learning and action adjustments in real-time."
        cp2.font.name = "Arial"
        cp2.font.size = Pt(16)
        cp2.font.color.rgb = color_text
    # Slide 10: 3D Digital Twin Screenshot
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    img_cesium_path = Path("docs/screenshots/cesium_3d_twin.jpg")
    if img_cesium_path.exists():
        slide.shapes.add_picture(str(img_cesium_path), Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    else:
        add_slide_header(slide, "Tokyo 3D Digital Twin", "CesiumJS Visualisation")
        card_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
        ctf = card_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = "🖥️ 3D Digital Twin Visualization"
        cp.font.name = "Arial"
        cp.font.size = Pt(22)
        cp.font.bold = True
        cp.font.color.rgb = color_primary
        cp.space_after = Pt(16)
        cp2 = ctf.add_paragraph()
        cp2.text = "Renders a high-fidelity 3D spatial view of Tokyo towers, modeling building geometry and NVIDIA Sionna ray-tracing coverage maps."
        cp2.font.name = "Arial"
        cp2.font.size = Pt(16)
        cp2.font.color.rgb = color_text

    # Slide 11: 3D Digital Twin Explanation
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Tokyo 3D Digital Twin Sandbox", "Spatial Network Twin")
    
    bullets_twin = [
        ("Urban Spatial Database", "Models 22 macro site towers across Shinjuku, Shibuya, and Minato-ku, capturing engineering azimuths, heights, mechanical/electrical tilts, and E2/A1 RIC interfaces."),
        ("NVIDIA Sionna Ray-Tracing", "Replaces Okumura-Hata formulas with GPU ray-tracing. Simulates coverage blocking, reflection, and diffraction across sub-6 GHz (n77/n78) and mmWave (n257) frequencies."),
        ("UE Mobility & Corridor Handovers", "Simulates real-world traffic paths (Yamanote Line trains, expressways) and visualizes serving sector paths (green: success, red: radio link failures)."),
        ("Operational Business ROI", "By using the twin to optimize CIO offsets using RL, handover failures drop from 2.5% to 0.5%, saving billions of yen in NOC triage and subscriber churn.")
    ]
    add_bullet_list(slide, bullets_twin, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
    ctf = card_box.text_frame
    ctf.word_wrap = True
    
    cp = ctf.paragraphs[0]
    cp.text = "🔮 Physical & Digital Twin"
    cp.font.name = "Arial"
    cp.font.size = Pt(22)
    cp.font.bold = True
    cp.font.color.rgb = color_primary
    cp.space_after = Pt(16)
    
    cp2 = ctf.add_paragraph()
    cp2.text = "The 3D twin models actual building geometry and radio physics. It allows the RL agent to test and validate network optimization policies in a safe, high-fidelity sandbox before deploying to live baseband nodes."
    cp2.font.name = "Arial"
    cp2.font.size = Pt(16)
    cp2.font.color.rgb = color_text

    # Slide 12: Results & ONNX
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "100k Sweep & ONNX Export", "Phase 3 Results")

    
    bullets_p5 = [
        ("100k Sweeps (16/16 Completed)", "PPO trained across 4 variants x 4 scenarios. Variant v2 baseline achieves HOSR: 99.99% and Ping-Pong: 0.00%."),
        ("Optuna Hyperparameter Tuning", "Optimized learning rate, batch size, rollout steps, and discount factor to resolve baseline HOSR stagnation."),
        ("ONNX Exporter (Complete)", "export_onnx.py converts PyTorch model zip checkpoint into deployment-ready ONNX models for CU-CP K8s deployment."),
        ("Tensor Clamping Bounds", "Applied torch.clamp matching environment action limits [-2.0, 2.0] to guarantee ONNX predictions match SB3 deterministic predictions.")
    ]
    add_bullet_list(slide, bullets_p5, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))
    
    card_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(5.0))
    ctf = card_box.text_frame
    ctf.word_wrap = True
    
    cp = ctf.paragraphs[0]
    cp.text = "🎯 Verification Performance"
    cp.font.name = "Arial"
    cp.font.size = Pt(22)
    cp.font.bold = True
    cp.font.color.rgb = color_primary
    cp.space_after = Pt(16)
    
    stats = [
        ("Random Baseline HOSR", "79.25%"),
        ("Trained Agent HOSR", "99.99%"),
        ("Absolute HOSR Delta", "+20.74%"),
        ("Trained Agent PP Rate", "0.00%"),
        ("Regression Tests", "262/262 passed")
    ]
    for lbl, val in stats:
        sp_stat = ctf.add_paragraph()
        sp_stat.text = f"• {lbl}: {val}"
        sp_stat.font.name = "Arial"
        sp_stat.font.size = Pt(16)
        sp_stat.font.bold = True
        sp_stat.font.color.rgb = color_text
        sp_stat.space_after = Pt(4)

    # Save the presentation
    output_path = Path("docs/presentation.pptx")
    prs.save(str(output_path))
    print(f"Presentation saved successfully to {output_path}!")


if __name__ == "__main__":
    create_presentation()
